"""Document parser factory — supports multiple file formats.

M1-T9: Parsers for PDF, Word, Excel, PPT, Markdown, and plain text.
Returns a unified Document model with metadata.
"""

from __future__ import annotations

import logging
import mimetypes
from abc import ABC, abstractmethod
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

from pydantic import BaseModel, Field

logger = logging.getLogger("fde.rag.parser")

# ══════════════════════════════════════════════════════════════════
# Unified Document Model
# ══════════════════════════════════════════════════════════════════


class Document(BaseModel):
    """Unified document model returned by all parsers."""

    id: str = Field(default="", description="Document ID (auto-generated if empty)")
    content: str = Field(description="Extracted text content")
    metadata: dict[str, Any] = Field(default_factory=dict)
    source: str = Field(default="", description="Original file path or URL")
    mime_type: str = Field(default="text/plain", description="MIME type of the source")
    page_number: int | None = Field(default=None, description="Page number for multi-page docs")
    chunk_index: int | None = Field(default=None, description="Chunk index for chunked docs")


# ══════════════════════════════════════════════════════════════════
# Parser Exceptions
# ══════════════════════════════════════════════════════════════════


class ParserError(Exception):
    """Base exception for document parsing."""


class UnsupportedFormatError(ParserError):
    """Raised when no parser is available for the given format."""


# ══════════════════════════════════════════════════════════════════
# Abstract Base Parser
# ══════════════════════════════════════════════════════════════════


class BaseParser(ABC):
    """Abstract base for document parsers."""

    supported_extensions: set[str] = set()
    supported_mime_types: set[str] = set()

    @abstractmethod
    def parse(self, file_path: str, **kwargs: Any) -> list[Document]:
        """Parse a document and return list of Document objects.

        Args:
            file_path: Path to the file to parse.
            **kwargs: Additional parser-specific options.

        Returns:
            List of Document objects (one per page for multi-page docs).
        """

    @abstractmethod
    async def async_parse(self, file_path: str, **kwargs: Any) -> list[Document]:
        """Async parse a document."""

    def validate_path(self, file_path: str) -> None:
        """Validate that the file exists and is readable."""
        path = Path(file_path)
        if not path.exists():
            raise ParserError(f"File not found: {file_path}")
        if not path.is_file():
            raise ParserError(f"Not a file: {file_path}")


# ══════════════════════════════════════════════════════════════════
# Format Detection
# ══════════════════════════════════════════════════════════════════


_MIME_EXT_MAP: dict[str, str] = {
    "application/pdf": ".pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": ".docx",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": ".xlsx",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": ".pptx",
    "text/markdown": ".md",
    "text/plain": ".txt",
}


def detect_format(file_path: str) -> str:
    """Detect MIME type from file extension or content."""
    path = Path(file_path)
    ext = path.suffix.lower()

    # Direct extension mapping
    ext_to_mime = {
        ".pdf": "application/pdf",
        ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ".md": "text/markdown",
        ".markdown": "text/markdown",
        ".txt": "text/plain",
        ".csv": "text/plain",
        ".json": "text/plain",
        ".yaml": "text/plain",
        ".yml": "text/plain",
        ".xml": "text/plain",
        ".html": "text/plain",
        ".htm": "text/plain",
    }

    mime = ext_to_mime.get(ext)
    if mime:
        return mime

    # Fallback to mimetypes
    guessed, _ = mimetypes.guess_type(file_path)
    return guessed or "text/plain"


# ══════════════════════════════════════════════════════════════════
# Individual Parsers
# ══════════════════════════════════════════════════════════════════


class PdfParser(BaseParser):
    """Parser for PDF documents using PyMuPDF."""

    supported_extensions = {".pdf"}
    supported_mime_types = {"application/pdf"}

    def parse(self, file_path: str, **kwargs: Any) -> list[Document]:
        self.validate_path(file_path)

        try:
            import pymupdf
        except ImportError:
            raise ParserError("PyMuPDF not installed. Install: pip install pymupdf")

        docs: list[Document] = []
        path = Path(file_path)

        with pymupdf.open(file_path) as pdf:
            for page_num, page in enumerate(pdf, start=1):
                text = page.get_text().strip()
                if not text:
                    continue
                docs.append(
                    Document(
                        content=text,
                        metadata={
                            "file_name": path.name,
                            "page": page_num,
                            "total_pages": len(pdf),
                            "format": "pdf",
                        },
                        source=str(path),
                        mime_type="application/pdf",
                        page_number=page_num,
                    )
                )

        logger.info("Parsed PDF '%s': %d pages", path.name, len(docs))
        return docs

    async def async_parse(self, file_path: str, **kwargs: Any) -> list[Document]:
        """Async parse — PyMuPDF is sync, runs in thread pool."""
        import asyncio

        return await asyncio.to_thread(self.parse, file_path, **kwargs)


class DocxParser(BaseParser):
    """Parser for Word (.docx) documents using python-docx."""

    supported_extensions = {".docx"}
    supported_mime_types = {
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    }

    def parse(self, file_path: str, **kwargs: Any) -> list[Document]:
        self.validate_path(file_path)

        try:
            from docx import Document as DocxDocument
        except ImportError:
            raise ParserError("python-docx not installed. Install: pip install python-docx")

        path = Path(file_path)
        doc = DocxDocument(file_path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        content = "\n".join(paragraphs)

        if not content:
            logger.warning("Docx '%s' has no extractable text", path.name)
            return []

        return [
            Document(
                content=content,
                metadata={
                    "file_name": path.name,
                    "paragraphs": len(paragraphs),
                    "format": "docx",
                },
                source=str(path),
                mime_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            )
        ]

    async def async_parse(self, file_path: str, **kwargs: Any) -> list[Document]:
        import asyncio

        return await asyncio.to_thread(self.parse, file_path, **kwargs)


class XlsxParser(BaseParser):
    """Parser for Excel (.xlsx) documents using openpyxl."""

    supported_extensions = {".xlsx"}
    supported_mime_types = {
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    }

    def parse(self, file_path: str, **kwargs: Any) -> list[Document]:
        self.validate_path(file_path)

        try:
            import openpyxl
        except ImportError:
            raise ParserError("openpyxl not installed. Install: pip install openpyxl")

        path = Path(file_path)
        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        docs: list[Document] = []

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows_text: list[str] = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                line = "\t".join(cells).strip()
                if line:
                    rows_text.append(line)

            content = "\n".join(rows_text)
            if content:
                docs.append(
                    Document(
                        content=content,
                        metadata={
                            "file_name": path.name,
                            "sheet": sheet_name,
                            "rows": len(rows_text),
                            "format": "xlsx",
                        },
                        source=str(path),
                        mime_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    )
                )

        wb.close()
        logger.info("Parsed XLSX '%s': %d sheets", path.name, len(docs))
        return docs

    async def async_parse(self, file_path: str, **kwargs: Any) -> list[Document]:
        import asyncio

        return await asyncio.to_thread(self.parse, file_path, **kwargs)


class PptxParser(BaseParser):
    """Parser for PowerPoint (.pptx) documents using python-pptx."""

    supported_extensions = {".pptx"}
    supported_mime_types = {
        "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    }

    def parse(self, file_path: str, **kwargs: Any) -> list[Document]:
        self.validate_path(file_path)

        try:
            from pptx import Presentation
        except ImportError:
            raise ParserError("python-pptx not installed. Install: pip install python-pptx")

        path = Path(file_path)
        prs = Presentation(file_path)
        docs: list[Document] = []

        for slide_idx, slide in enumerate(prs.slides, start=1):
            texts: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text)
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_texts = [cell.text.strip() for cell in row.cells]
                        texts.append(" | ".join(row_texts))

            content = "\n".join(texts)
            if content:
                docs.append(
                    Document(
                        content=content,
                        metadata={
                            "file_name": path.name,
                            "slide": slide_idx,
                            "total_slides": len(prs.slides),
                            "format": "pptx",
                        },
                        source=str(path),
                        mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        page_number=slide_idx,
                    )
                )

        logger.info("Parsed PPTX '%s': %d slides", path.name, len(docs))
        return docs

    async def async_parse(self, file_path: str, **kwargs: Any) -> list[Document]:
        import asyncio

        return await asyncio.to_thread(self.parse, file_path, **kwargs)


class MarkdownParser(BaseParser):
    """Parser for Markdown (.md) documents."""

    supported_extensions = {".md", ".markdown"}
    supported_mime_types = {"text/markdown"}

    def parse(self, file_path: str, **kwargs: Any) -> list[Document]:
        self.validate_path(file_path)
        path = Path(file_path)
        content = path.read_text(encoding="utf-8")

        if not content.strip():
            logger.warning("Markdown file '%s' is empty", path.name)
            return []

        return [
            Document(
                content=content,
                metadata={
                    "file_name": path.name,
                    "format": "markdown",
                    "char_count": len(content),
                    "line_count": content.count("\n") + 1,
                },
                source=str(path),
                mime_type="text/markdown",
            )
        ]

    async def async_parse(self, file_path: str, **kwargs: Any) -> list[Document]:
        import asyncio

        return await asyncio.to_thread(self.parse, file_path, **kwargs)


class PlainTextParser(BaseParser):
    """Parser for plain text files (.txt, .csv, .json, etc.)."""

    supported_extensions = {".txt", ".csv", ".json", ".yaml", ".yml", ".xml", ".html", ".htm"}
    supported_mime_types = {"text/plain"}

    def parse(self, file_path: str, **kwargs: Any) -> list[Document]:
        self.validate_path(file_path)
        path = Path(file_path)

        try:
            content = path.read_text(encoding="utf-8")
        except UnicodeDecodeError:
            # Fallback for non-UTF-8 files
            content = path.read_text(encoding="latin-1")

        if not content.strip():
            logger.warning("Text file '%s' is empty", path.name)
            return []

        return [
            Document(
                content=content,
                metadata={
                    "file_name": path.name,
                    "format": "text",
                    "char_count": len(content),
                    "line_count": content.count("\n") + 1,
                },
                source=str(path),
                mime_type="text/plain",
            )
        ]

    async def async_parse(self, file_path: str, **kwargs: Any) -> list[Document]:
        import asyncio

        return await asyncio.to_thread(self.parse, file_path, **kwargs)


# ══════════════════════════════════════════════════════════════════
# Parser Factory
# ══════════════════════════════════════════════════════════════════


class ParserFactory:
    """Factory for creating document parsers based on file format.

    Usage:
        factory = ParserFactory()
        parser = factory.get_parser("report.pdf")
        docs = parser.parse("report.pdf")
    """

    def __init__(self) -> None:
        self._parsers: list[BaseParser] = self._register_defaults()

    @staticmethod
    def _register_defaults() -> list[BaseParser]:
        """Register all built-in parsers."""
        return [
            PdfParser(),
            DocxParser(),
            XlsxParser(),
            PptxParser(),
            MarkdownParser(),
            PlainTextParser(),
        ]

    def get_parser(self, file_path: str) -> BaseParser:
        """Get the appropriate parser for a file.

        Args:
            file_path: Path to the file.

        Returns:
            A parser instance matching the file format.

        Raises:
            UnsupportedFormatError: If no parser supports the format.
        """
        mime_type = detect_format(file_path)
        ext = Path(file_path).suffix.lower()

        # Try MIME type match first, then extension match
        for parser in self._parsers:
            if mime_type in parser.supported_mime_types:
                return parser
            if ext in parser.supported_extensions:
                return parser

        raise UnsupportedFormatError(
            f"No parser available for '{file_path}' (MIME: {mime_type}, ext: {ext})"
        )

    def parse(self, file_path: str, **kwargs: Any) -> list[Document]:
        """Parse a document by auto-detecting format.

        Convenience method that combines get_parser + parse.

        Args:
            file_path: Path to the file.
            **kwargs: Parser-specific options.

        Returns:
            List of Document objects.
        """
        parser = self.get_parser(file_path)
        return parser.parse(file_path, **kwargs)

    async def async_parse(self, file_path: str, **kwargs: Any) -> list[Document]:
        """Async parse a document by auto-detecting format."""
        parser = self.get_parser(file_path)
        return await parser.async_parse(file_path, **kwargs)

    def list_supported_formats(self) -> list[dict[str, Any]]:
        """List all supported formats with their extensions."""
        formats: list[dict[str, Any]] = []
        for parser in self._parsers:
            formats.append(
                {
                    "parser": parser.__class__.__name__,
                    "extensions": sorted(parser.supported_extensions),
                    "mime_types": sorted(parser.supported_mime_types),
                }
            )
        return formats