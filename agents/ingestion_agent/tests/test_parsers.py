"""P2b 解析器测试：xlsx/csv/docx/pptx/pdf 多格式解析为 ParsedDocument。"""

from __future__ import annotations

import io

import openpyxl

from agents.ingestion_agent.parsers import (
    BlockType,
    ParsedDocument,
    detect_file_type,
    parse_file,
)


def _xlsx_bytes() -> bytes:
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.append(["客户名称", "Order-No", "城市"])
    ws.append(["阿里巴巴", "SO-1", "杭州"])
    ws.append(["腾讯", "SO-2", "深圳"])
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _csv_bytes() -> bytes:
    return "客户名称,Order-No,城市\n阿里巴巴,SO-1,杭州\n腾讯,SO-2,深圳\n".encode()


def _docx_bytes() -> bytes:
    import docx

    d = docx.Document()
    d.add_heading("合同清单", level=1)
    d.add_paragraph("阿里巴巴 总部位于杭州。")
    d.add_paragraph("腾讯科技 总部位于深圳。")
    t = d.add_table(rows=1, cols=2)
    t.rows[0].cells[0].text = "客户"
    t.rows[0].cells[1].text = "城市"
    r = t.add_row().cells
    r[0].text = "阿里巴巴"
    r[1].text = "杭州"
    buf = io.BytesIO()
    d.save(buf)
    return buf.getvalue()


def _pptx_bytes() -> bytes:
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "季度复盘"
    tb = slide.shapes.add_textbox(0, 0, 100, 100)
    tb.text_frame.text = "本季度签约 3 家客户。"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _pdf_bytes() -> bytes:
    import fitz

    doc = fitz.open()
    page = doc.new_page()
    # 注：fitz.insert_text 用 base-14 字体不含 CJK 字形，故测试用 ASCII 文本；
    # 真实 PDF（Word/LaTeX 导出）内嵌 CJK 字体时 pdfplumber 可正常抽取中文。
    page.insert_text((50, 50), "Alibaba headquarters in Hangzhou.")
    page.insert_text((50, 70), "Tencent headquarters in Shenzhen.")
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def test_detect_by_extension() -> None:
    assert detect_file_type("a.xlsx", b"") == "excel"
    assert detect_file_type("a.csv", b"") == "csv"
    assert detect_file_type("a.pdf", b"") == "pdf"
    assert detect_file_type("a.docx", b"") == "docx"
    assert detect_file_type("a.pptx", b"") == "pptx"
    assert detect_file_type("a.bin", b"") == "unknown"


def test_parse_excel_returns_table_block() -> None:
    doc = parse_file("sample.xlsx", _xlsx_bytes())
    assert isinstance(doc, ParsedDocument)
    assert doc.meta["file_type"] == "excel"
    assert len(doc.blocks) == 1
    b = doc.blocks[0]
    assert b.kind == BlockType.TABLE
    assert "客户名称" in b.table_headers
    assert len(b.table) == 2
    assert b.table[0]["客户名称"] == "阿里巴巴"


def test_parse_csv_returns_table_block() -> None:
    doc = parse_file("sample.csv", _csv_bytes())
    assert doc.meta["file_type"] == "csv"
    b = doc.blocks[0]
    assert b.kind == BlockType.TABLE
    assert b.table[0]["客户名称"] == "阿里巴巴"
    assert b.table[1]["城市"] == "深圳"


def test_parse_docx_text_and_table() -> None:
    doc = parse_file("sample.docx", _docx_bytes())
    assert doc.meta["file_type"] == "docx"
    kinds = [b.kind for b in doc.blocks]
    assert BlockType.HEADING in kinds
    assert BlockType.TEXT in kinds
    assert BlockType.TABLE in kinds
    tbl = next(b for b in doc.blocks if b.kind == BlockType.TABLE)
    assert tbl.table[0]["客户"] == "阿里巴巴"


def test_parse_pptx_text_block() -> None:
    doc = parse_file("sample.pptx", _pptx_bytes())
    assert doc.meta["file_type"] == "pptx"
    texts = [b.text for b in doc.blocks if b.kind == BlockType.TEXT]
    assert any("签约" in t for t in texts)


def test_parse_pdf_text_block() -> None:
    doc = parse_file("sample.pdf", _pdf_bytes())
    assert doc.meta["file_type"] == "pdf"
    all_text = " ".join(b.text for b in doc.blocks if b.kind == BlockType.TEXT)
    assert "Hangzhou" in all_text
