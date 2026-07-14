"""多格式文件解析器（P2b / 完整本地文件入库）。

把 ``.xlsx / .xlsm / .csv / .pdf / .docx / .pptx`` 解析为统一的 ``ParsedDocument``
（blocks：``text`` / ``table`` / ``heading``），供 normalization + chunking + 入库复用。

解析库均为**惰性导入**（函数内 ``import``），缺失某个库只影响对应类型，不会让整个
ingestion router 注册失败。规划允许 Docling 经 P0.5 spike 不通过时回退到
``pdfplumber + python-docx + openpyxl``——本实现即采用该回退路径，避免 ARM 机器上
约 1.5G 的 torch 依赖（见 docs/master-delivery-plan.md P0.5 / P2b）。

**P2-A 升级（Docling 可选后端）**：PDF 解析支持 ``docling`` 后端（``DocumentConverter``
输出 Markdown → 结构化 Block），通过环境变量 ``FDE_PARSER_BACKEND`` 控制：
``auto``（docling 优先、未安装则回退 pdfplumber，默认）/ ``docling``（强制，未装报错）/
``pdfplumber``（强制回退）。Docling 依赖 torch VLM，列为可选依赖
（``requirements-docling.txt``），不进入主依赖树，服务器保持零 torch。
"""

from __future__ import annotations

import csv
import io
import os
import re
from dataclasses import dataclass, field
from enum import StrEnum
from pathlib import Path
from typing import Any

# PDF 解析后端选择（P2-A）。docling 为可选增强，pdfplumber 为默认回退。
_PDF_BACKEND = os.getenv("FDE_PARSER_BACKEND", "auto").lower()

# 扩展名 → 解析器类型
_EXT_LOADERS: dict[str, str] = {
    ".xlsx": "excel",
    ".xlsm": "excel",
    ".csv": "csv",
    ".pdf": "pdf",
    ".docx": "docx",
    ".pptx": "pptx",
}


class BlockType(StrEnum):
    """解析块类型。"""

    TEXT = "text"
    TABLE = "table"
    HEADING = "heading"


@dataclass
class Block:
    """一个解析块（文本段 / 表格 / 标题）。

    ``table`` 与 ``table_headers`` 仅当 ``kind == TABLE`` 时使用；行字典的键为
    **原始表头**（保留空格），归一化在 normalization 层完成，避免与映射契约错位。
    """

    kind: BlockType
    text: str = ""
    table: list[dict[str, Any]] | None = None
    table_headers: list[str] | None = None
    loc: dict[str, Any] = field(default_factory=dict)


@dataclass
class ParsedDocument:
    """一次解析的产物。"""

    doc_type: str
    source_ref: str
    blocks: list[Block] = field(default_factory=list)
    meta: dict[str, Any] = field(default_factory=dict)


def detect_file_type(filename: str, data: bytes) -> str:
    """按扩展名（辅以 magic bytes 兜底）判断文件类型。"""
    ext = Path(filename).suffix.lower()
    if ext in _EXT_LOADERS:
        return _EXT_LOADERS[ext]
    if data[:4] == b"%PDF":
        return "pdf"
    if data[:2] == b"PK":
        return "docx"  # docx/pptx 均为 zip，扩展名兜底更准确
    return "unknown"


def _grid_to_table(grid: list[list[Any]]) -> tuple[list[str], list[dict[str, Any]]]:
    """把 PDF/DOCX/PPTX 的二维单元格网格转成 (原始表头, 行字典列表)。

    行字典键使用**原始表头**（保留空格、自动去重），归一化交给 normalization 层。
    """
    if not grid:
        return [], []
    raw_headers = [str(c).strip() if c is not None else "" for c in grid[0]]
    headers: list[str] = []
    seen: dict[str, int] = {}
    for i, h in enumerate(raw_headers):
        base = h if h else f"col_{i}"
        if base in seen:
            seen[base] += 1
            base = f"{base}_{seen[base]}"
        else:
            seen[base] = 0
        headers.append(base)
    rows: list[dict[str, Any]] = []
    for r in grid[1:]:
        if r is None or all(c is None or str(c).strip() == "" for c in r):
            continue
        row = {headers[i]: (r[i] if i < len(r) else None) for i in range(len(headers))}
        rows.append(row)
    return headers, rows


def _table_block(headers: list[str], rows: list[dict[str, Any]], loc: dict[str, Any]) -> Block:
    return Block(kind=BlockType.TABLE, table=rows, table_headers=headers, loc=loc)


def parse_file(filename: str, data: bytes) -> ParsedDocument:
    """解析任意受支持的文件，返回 ``ParsedDocument``。"""
    ftype = detect_file_type(filename, data)
    if ftype == "excel":
        return _parse_excel(data, filename)
    if ftype == "csv":
        return _parse_csv(data, filename)
    if ftype == "pdf":
        return _parse_pdf(data, filename)
    if ftype == "docx":
        return _parse_docx(data, filename)
    if ftype == "pptx":
        return _parse_pptx(data, filename)
    raise ValueError(f"不支持的文件类型: {filename}")


def _parse_excel(data: bytes, filename: str) -> ParsedDocument:
    try:
        import openpyxl
    except ImportError as exc:  # pragma: no cover - openpyxl 为声明依赖
        raise RuntimeError("openpyxl 未安装，无法解析 Excel") from exc
    wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    blocks: list[Block] = []
    for ws in wb.worksheets:
        grid = list(ws.iter_rows(values_only=True))
        if not grid:
            continue
        headers, rows = _grid_to_table(grid)
        if rows:
            blocks.append(_table_block(headers, rows, {"sheet": ws.title}))
    return ParsedDocument(
        doc_type="excel_upload",
        source_ref=filename,
        blocks=blocks,
        meta={
            "content_type": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            "file_type": "excel",
        },
    )


def _parse_csv(data: bytes, filename: str) -> ParsedDocument:
    text = data.decode("utf-8-sig", errors="replace")
    reader = csv.reader(io.StringIO(text))
    grid = list(reader)
    headers, rows = _grid_to_table(grid)
    if rows:
        blocks = [_table_block(headers, rows, {"sheet": "csv"})]
    else:
        blocks = []
    return ParsedDocument(
        doc_type="csv_upload",
        source_ref=filename,
        blocks=blocks,
        meta={"content_type": "text/csv", "file_type": "csv"},
    )


def _parse_pdf(data: bytes, filename: str) -> ParsedDocument:
    """PDF 解析：后端由 ``FDE_PARSER_BACKEND`` 决定（P2-A）。

    - ``pdfplumber``：强制使用轻量回退解析器。
    - ``docling``：强制使用 Docling（未安装则抛 ImportError）。
    - ``auto``（默认）：Docling 优先，未安装时自动回退 pdfplumber。
    """
    if _PDF_BACKEND == "pdfplumber":
        return _parse_pdf_pdfplumber(data, filename)
    if _PDF_BACKEND == "docling":
        return _parse_pdf_docling(data, filename)
    # auto：docling 优先，失败回退 pdfplumber
    try:
        return _parse_pdf_docling(data, filename)
    except ImportError:
        return _parse_pdf_pdfplumber(data, filename)


def _parse_pdf_pdfplumber(data: bytes, filename: str) -> ParsedDocument:
    try:
        import pdfplumber
    except ImportError as exc:  # pragma: no cover - pdfplumber 为声明依赖
        raise RuntimeError("pdfplumber 未安装，无法解析 PDF") from exc
    blocks: list[Block] = []
    with pdfplumber.open(io.BytesIO(data)) as pdf:
        for pi, page in enumerate(pdf.pages, start=1):
            text = page.extract_text() or ""
            for bi, para in enumerate(p for p in text.split("\n\n") if p.strip()):
                blocks.append(
                    Block(kind=BlockType.TEXT, text=para.strip(), loc={"page": pi, "block": bi})
                )
            for ti, tbl in enumerate(page.extract_tables() or []):
                headers, rows = _grid_to_table(tbl)
                if rows:
                    blocks.append(_table_block(headers, rows, {"page": pi, "table": ti}))
    return ParsedDocument(
        doc_type="pdf_upload",
        source_ref=filename,
        blocks=blocks,
        meta={"content_type": "application/pdf", "file_type": "pdf", "parser": "pdfplumber"},
    )


def _parse_markdown_to_blocks(md: str, *, source: str = "docling") -> list[Block]:
    """将 Markdown 文本解析为结构化 Block 列表（P2-A Docling 复用）。

    表格（GitHub 风格 ```| a | b |``` 块）转为 TABLE block，其余按标题层级
    标记为 HEADING / TEXT。表格解析复用 ``_grid_to_table``，与现有管线一致。
    """
    lines = md.split("\n")
    blocks: list[Block] = []
    i = 0
    n = len(lines)
    table_buf: list[str] = []

    def _flush_table() -> None:
        nonlocal table_buf
        if len(table_buf) < 2:
            # 不足表头+分隔行，当作文本处理
            text = "\n".join(table_buf).strip()
            if text:
                blocks.append(Block(kind=BlockType.TEXT, text=text, loc={"source": source}))
            table_buf = []
            return
        # 识别分隔行（|---|---|）
        sep_idx = -1
        for idx, row in enumerate(table_buf):
            cells = [c.strip() for c in row.strip().strip("|").split("|")]
            if idx > 0 and all(re.fullmatch(r":?-{2,}:?", c) for c in cells if c != ""):
                sep_idx = idx
                break
        if sep_idx <= 0:
            text = "\n".join(table_buf).strip()
            if text:
                blocks.append(Block(kind=BlockType.TEXT, text=text, loc={"source": source}))
            table_buf = []
            return
        # 切掉分隔行，仅保留表头 + 数据行
        grid = []
        for row in table_buf[:sep_idx] + table_buf[sep_idx + 1 :]:
            grid.append([c.strip() for c in row.strip().strip("|").split("|")])
        headers, rows = _grid_to_table(grid)
        if rows:
            blocks.append(_table_block(headers, rows, {"source": source}))
        table_buf = []

    while i < n:
        line = lines[i]
        stripped = line.strip()
        # 表格行检测：以 | 包裹或起始/结束含 |
        if stripped.startswith("|") or (stripped.endswith("|") and "|" in stripped[:-1]):
            table_buf.append(line)
            i += 1
            continue
        # 非表格行：先 flush 累积的表格
        if table_buf:
            _flush_table()
        if not stripped:
            i += 1
            continue
        # 标题（# 层级）
        m = re.match(r"^(#{1,6})\s+(.*)$", stripped)
        if m:
            blocks.append(
                Block(kind=BlockType.HEADING, text=m.group(2).strip(), loc={"source": source})
            )
        else:
            blocks.append(Block(kind=BlockType.TEXT, text=stripped, loc={"source": source}))
        i += 1
    if table_buf:
        _flush_table()
    return blocks


def _parse_pdf_docling(data: bytes, filename: str) -> ParsedDocument:
    """Docling 后端（P2-A）：``DocumentConverter`` → Markdown → 结构化 Block。

    惰性导入 docling；未安装时抛 ``ImportError``，由 ``_parse_pdf`` 回退。
    表格/布局检测依赖 torch VLM，属可选重型依赖（见模块 docstring）。
    """
    try:
        from docling.datamodel.base_models import DocumentStream
        from docling.document_converter import DocumentConverter
    except ImportError as exc:  # pragma: no cover - docling 为可选依赖
        raise ImportError("docling 未安装，无法使用 docling 后端") from exc

    stream = DocumentStream(name=filename, stream=io.BytesIO(data))
    try:
        result = DocumentConverter().convert(stream)
    except Exception as exc:  # 转换失败（损坏/扫描件）→ 让上层回退
        raise RuntimeError(f"docling 转换失败: {exc}") from exc
    md = result.document.export_to_markdown()
    blocks = _parse_markdown_to_blocks(md, source="docling")
    return ParsedDocument(
        doc_type="pdf_upload",
        source_ref=filename,
        blocks=blocks,
        meta={"content_type": "application/pdf", "file_type": "pdf", "parser": "docling"},
    )


def parse_pdf_backend() -> str:
    """返回当前生效的 PDF 解析后端（解析 ``auto`` 为实际后端）。

    用于路由层/健康自检报告当前解析能力，以及测试断言后端选择。
    """
    if _PDF_BACKEND == "pdfplumber":
        return "pdfplumber"
    if _PDF_BACKEND == "docling":
        return "docling"
    # auto：探测 docling 是否可导入
    try:
        import docling  # noqa: F401

        return "docling"
    except ImportError:
        return "pdfplumber"


def _parse_docx(data: bytes, filename: str) -> ParsedDocument:
    try:
        import docx
    except ImportError as exc:  # pragma: no cover - python-docx 为声明依赖
        raise RuntimeError("python-docx 未安装，无法解析 DOCX") from exc
    document = docx.Document(io.BytesIO(data))
    blocks: list[Block] = []
    for pi, para in enumerate(document.paragraphs):
        txt = para.text.strip()
        if not txt:
            continue
        style = (para.style.name or "") if para.style else ""
        kind = BlockType.HEADING if style.startswith("Heading") else BlockType.TEXT
        blocks.append(Block(kind=kind, text=txt, loc={"paragraph": pi}))
    for ti, table in enumerate(document.tables):
        grid = [[c.text for c in row.cells] for row in table.rows]
        headers, rows = _grid_to_table(grid)
        if rows:
            blocks.append(_table_block(headers, rows, {"table": ti}))
    return ParsedDocument(
        doc_type="docx_upload",
        source_ref=filename,
        blocks=blocks,
        meta={
            "content_type": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "file_type": "docx",
        },
    )


def _parse_pptx(data: bytes, filename: str) -> ParsedDocument:
    try:
        from pptx import Presentation
    except ImportError as exc:  # pragma: no cover - python-pptx 为声明依赖
        raise RuntimeError("python-pptx 未安装，无法解析 PPTX") from exc
    prs = Presentation(io.BytesIO(data))
    blocks: list[Block] = []
    for si, slide in enumerate(prs.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    texts.append(t)
            if shape.has_table:
                grid = [[c.text for c in row.cells] for row in shape.table.rows]
                headers, rows = _grid_to_table(grid)
                if rows:
                    blocks.append(_table_block(headers, rows, {"slide": si, "table": 0}))
        if texts:
            blocks.append(Block(kind=BlockType.TEXT, text="\n".join(texts), loc={"slide": si}))
    return ParsedDocument(
        doc_type="pptx_upload",
        source_ref=filename,
        blocks=blocks,
        meta={
            "content_type": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "file_type": "pptx",
        },
    )


__all__ = [
    "Block",
    "BlockType",
    "ParsedDocument",
    "detect_file_type",
    "parse_file",
    "parse_pdf_backend",
]
